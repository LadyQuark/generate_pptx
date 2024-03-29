import traceback

from pathlib import Path
from pptx import Presentation

from common import create_text_chunks, find_and_replace_diagrams, find_and_replace_OLE_photos, find_and_replace_OLE, print_shape_type
from utils import copy_shapes, remove_shape, _object_rels

class PresentationManager(object):
    """Contains Presentation object and functions to manage it"""
    
    # Character limit for content text in single slide
    MAX_CONTENT_LIMIT=2250

    def __init__(self, file_path=None, template_slide_index=1, slide_size=()):
        # Since presentation.Presentation class not intended to be constructed directly, 
        # using pptx.Presentation() to open presentation
        if file_path and Path(file_path).exists():
            self.presentation = Presentation(file_path)
            print("Loaded presentation from:", file_path)
        else:
            if file_path:
                print(f"Could not load {file_path}")  
            self.presentation = Presentation()
            print("New presentation object loaded")

        if slide_size:
            height, width = slide_size
            self.set_slide_size(height, width)

        # Setting index of slide to be used as a template
        self.template_slide_index = template_slide_index

        # Get index of Blank slide layout
        layout_items_count = [len(layout.placeholders) for layout in self.presentation.slide_layouts]
        min_items = min(layout_items_count)
        self.blank_layout_id = layout_items_count.index(min_items)

        for slide in self.presentation.slides:
            find_and_replace_diagrams(slide)
            find_and_replace_OLE_photos(slide)
            # find_and_replace_OLE(slide)

    @property
    def xml_slides(self):
        return self.presentation.slides._sldIdLst

    @property
    def _blank_slide_layout(self):        
        return self.presentation.slide_layouts[self.blank_layout_id]
    
    @property
    def total_slides(self):
        return len(self.presentation.slides)
    
    def set_slide_size(self, height, width):
        self.presentation.slide_height = height
        self.presentation.slide_width = width        

    def duplicate_slide(self, index, destination=None):
        """
        Duplicates the slide with the given index. Adds slide to the end of the presentation
        """
        source = self.presentation.slides[index]
        destination = destination or self
        # Adds blank slide to end
        blank_slide_layout = destination._blank_slide_layout
        dest = destination.presentation.slides.add_slide(blank_slide_layout)

        # Remove all shapes from the default layout
        for shape in dest.shapes:
            remove_shape(shape)

        # Copy all existing shapes
        copy_shapes(source.shapes, dest)

        # Copy existing references of known type
        # e.g. hyperlinks
        known_refs = [
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject"
        ]
        for rel in _object_rels(source.part):
            if rel.reltype in known_refs:
                if rel.is_external:
                    dest.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
                else:
                    dest.part.rels.get_or_add(rel.reltype, rel._target)

        # Copy all existing shapes
        if source.has_notes_slide:
            txt = source.notes_slide.notes_text_frame.text
            dest.notes_slide.notes_text_frame.text = txt

        return dest

    def move_slide(self, old_index, new_index):
        slides = list(self.xml_slides)
        self.xml_slides.remove(slides[old_index])
        self.xml_slides.insert(new_index, slides[old_index])

    def remove_slide(self, index):
        slides = list(self.xml_slides)
        self.xml_slides.remove(slides[index]) 

    def remove_all_slides(self):
        slides = list(self.xml_slides)
        for slide in slides:
            self.xml_slides.remove(slide)       
        

    def add_text_to_slide(self, index, text_content, title=""):
        """Adds title and content to slide at given index"""
        
        dest = self.presentation.slides[index]
        # Get title frame and content frame
        text_frames = []
        for shape in dest.shapes:
            if shape.has_text_frame:
                if 'Title' in shape.name:
                    title_frame = shape.text_frame
                else:
                    text_frames.append(shape.text_frame)
        # Choose first text frame as target
        content_frame = text_frames[0]
        
        # Clear content frame and add text
        content_frame.clear()
        p = content_frame.paragraphs[0]
        run = p.add_run()
        run.text = text_content

        # Clear title frame and add title
        title_frame.clear()
        p = title_frame.paragraphs[0]
        run = p.add_run()
        run.text = title

    def populate_slide(self, content, title=""):
        """Creates slides with given text and title, making more slides if text over limit"""
        
        duplicate_indices = []
        chunks = create_text_chunks(content, self.MAX_CONTENT_LIMIT)

        # Create slides for each chunk of text
        for chunk in chunks:
            slide_copy = self.duplicate_slide(self.template_slide_index)
            i = self.presentation.slides.index(slide_copy)
            duplicate_indices.append(i)
            self.add_text_to_slide(i, chunk, title)

        # Move all slides to just after template slide
        new_index = self.template_slide_index + 1
        for old_index in duplicate_indices:
            self.move_slide(old_index, new_index)
            new_index += 1


    def save(self, filepath, remove_template=False):
        """Saves presentation to given filepath and removes slide used as template"""

        if remove_template:
            print("Removing template", self.template_slide_index)
            self.remove_slide(self.template_slide_index)
        self.presentation.save(filepath)
        print("Saved presentation to:", filepath)


    @classmethod
    def copy_slide_to_other_presentation(cls, source, dest_filepath, slides_to_copy=[]):
        # Load presentation
        destination = PresentationManager(dest_filepath)
        # Copy presentation size if destination is empty
        if destination.total_slides == 0:
            height, width = source.presentation.slide_height, source.presentation.slide_width
            destination.set_slide_size(height, width)
        # If no slide numbers given, default to all slides
        if not slides_to_copy:
            slides_to_copy = range(source.total_slides)
        try:
            for i in slides_to_copy:
                source.duplicate_slide(i, destination)
            destination.save(dest_filepath)
            destination = Presentation(dest_filepath)
            destination.save(dest_filepath)
        except Exception:
            traceback.print_exc()   


    def _analyse_slide_elements(self, index, description=None):
        slide = self.presentation.slides[index]
        if description:
            print("*"*40, description, "*" * 40, sep="\n")
        for shape in slide.shapes:
            print_shape_type(shape)

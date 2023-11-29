from io import BytesIO
from lxml import etree
from lxml.etree import Element, SubElement
from pptx.oxml import parse_xml
from pptx.shapes.autoshape import Shape
from pptx.shapes.group import GroupShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pprint
pp = pprint.PrettyPrinter(indent=4)

NS = {'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram',
 'dsp': 'http://schemas.microsoft.com/office/drawing/2008/diagram',
 'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
 'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

def create_text_chunks(text, max_chunk_size=2250):
    chunks = []
    remaining = text
    while len(remaining) > max_chunk_size:
        # Determine best point `i` to truncate text where i < max_chunk_size
        i = max_chunk_size
        while (remaining[i] != "\n") and (remaining[i-1:i+1] != ". "):
            i -= 1
        # Split remaining text into two — append first to `chunks` and second to `remaining`
        chunks.append(remaining[:i])
        remaining = remaining[i:]
    # Append remaining text to `chunks`
    chunks.append(remaining)
    # Strip all chunks of trailing whitespace and newlines
    for i in range(len(chunks)):
        chunks[i] = chunks[i].strip("\n")
        chunks[i] = chunks[i].strip()
    
    return chunks


def find_and_replace_diagrams(slide):
    # Collect all diagrams in slide
    diagrams = []
    for shape in slide.shapes:
        if "Diagram" in shape.name:
            diagrams.append(shape)

    for diagram in diagrams:
        # Get matching source XML for diagram
        drawing_xml = get_drawing_xml(diagram)
        xfrm = []
        # xfrm = get_xfrm(diagram)
        # Remove diagram
        parent = diagram.element.getparent()
        parent.remove(diagram.element)
        # Get id number for next group shape
        # group_id = get_next_id(slide)
        next_id = slide.shapes._next_shape_id
        group_id = next_id
        # Make Shape objects
        new_shape_objects = shapes_from_drawing(drawing_xml, group_id + 1, parent)
        # Create new groupShape, attach shape objects, attach to slide
        add_group_to_slide(slide, group_id, new_shape_objects, xfrm)


def find_and_replace_OLE_photos(slide):
    shapes = slide.shapes
    # Collect all embedded OLE objects
    ole_photo_objs = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            if shape.ole_format.prog_id == "MSPhotoEd.3":
                ole_photo_objs.append(shape)

    for obj in ole_photo_objs:
        # Get picture object
        pic = obj.element.xpath(".//a:graphic/a:graphicData/*/*/p:oleObj/p:pic")[0]
        # Edit id and name of picture
        cNvPr = pic.xpath(".//p:nvPicPr/p:cNvPr")[0]
        # next_id = get_next_id(slide)
        next_id = slide.shapes._next_shape_id
        cNvPr.set("name", "Picture " + str(next_id))
        cNvPr.set("id", str(next_id))
        # Insert picture
        shapes._spTree.insert_element_before(pic, "p:extLst")
        # Remove OLE object
        parent = obj.element.getparent()
        parent.remove(obj.element)

def find_and_replace_OLE(slide):
    ole_photo_objs = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            if shape.ole_format.prog_id != "MSPhotoEd.3":
                ole_photo_objs.append(shape)

    for obj in ole_photo_objs:
        nsmap = obj.element.nsmap
        nsmap.update({"mc": "http://schemas.openxmlformats.org/markup-compatibility/2006"})
        left, top = obj.left, obj.top
        ole = obj.ole_format
        prog_id = ole.prog_id
        try:
            # ole_object = ole.element.xpath(".//*/*/p:oleObj")[0]
            # ole_object = next(ole.element.iter(".//mc:AlternateContent/mc:Fallback/p:oleObj"))
            a_blip = next(obj.element.iter("{%s}blip" % nsmap['a']))
            rel_id = next(value for _, value in a_blip.attrib.items())
        except StopIteration:
            ole_object = next(ole.element.iter("{%s}oleObj" % nsmap['p']))
            rel_id = ole_object.get("{%s}id" % ole_object.nsmap['r'])
        embed = obj.part.rels[rel_id]._target._blob

        parent = obj.element.getparent()
        parent.remove(obj.element)    
        
        slide.shapes.add_ole_object(
            object_file=BytesIO(embed),
            prog_id=prog_id,
            left=left,
            top=top,
        )


def get_next_id(slide):
    max_id = 1
    for shape in slide.shapes:
        if "Content Placeholder" not in shape.name:
            max_id = max(max_id, shape.shape_id)
    return max_id + 1

def get_drawing_xml(diagram):
    for _, rel in diagram.part.rels.items():
        if rel.target_partname.endswith("drawing1.xml"):
            drawing_xml = rel._target._blob    
            return drawing_xml   
def get_xfrm(diagram):
    xfrm = []
    for child in diagram.element.xpath("//p:xfrm"):
        print(child.tag)
        child.tag = "{%s}" % NS['a'] + "xfrm"
        xfrm.append(child) 
    return xfrm   


def shapes_from_drawing(drawing_xml, id_next, parent):
    new_shape_objects = []

    # Load drawing as XML element
    drawing = parse_xml(drawing_xml)
    nsmap = drawing.nsmap
    nsmap.update({"p":"http://schemas.openxmlformats.org/presentationml/2006/main"})
    etree.register_namespace("p", nsmap["p"])

    # Replace namespace `dsp` with `p`
    spTree = drawing[0]
    for el in drawing.iter("{%s}*" % nsmap['dsp']):
        tag = etree.QName(el)
        el.tag = etree.QName(nsmap['p'], tag.localname).text   

    # Change properties of all shapes (elements with p:sp tag)
    shapes = spTree.findall("p:sp", nsmap)
    len(shapes)
    for shape in shapes:
        cNvPr = shape.find(".//p:nvSpPr/p:cNvPr", nsmap)
        cNvPr.set("id", str(id_next))
        id_next += 1
        cNvPr.set("name", f"Freeform {id_next}")
        shape.set("has_ph_elm", "False")
        nvSpPr = shape.find(".//p:nvSpPr", nsmap)
        etree.SubElement(nvSpPr, etree.QName(nsmap['p'], "nvPr").text)     
        shape_obj = Shape(shape, parent)
        new_shape_objects.append(shape_obj)

    return new_shape_objects


def add_group_to_slide(slide, group_id, shapes, xfrm=[]):
    P = "{%s}" % NS['p']
    new_group = slide.shapes.add_group_shape(shapes=shapes)
    new_group.name = "Group " + str(group_id)

    nvGrpSpPr   = SubElement(new_group.element, P + "nvGrpSpPr")
    cNvPr       = SubElement(nvGrpSpPr, P + "cNvPr",
                                attrib={
                                    "id":   str(group_id),
                                    "name": "Group " + str(group_id),
                                })
    cNvGrpSpPr  = SubElement(nvGrpSpPr, P + "cNvGrpSpPr")
    nvPr        = SubElement(nvGrpSpPr, P + "nvPr")
    grpSpPr     = SubElement(new_group.element, P + "grpSpPr")  
    
    if xfrm:
        grpSpPr.extend(xfrm)


def print_shape_type(shape, indent=0):
    shape_string = "\t" * indent + "{:<22s} | {:<25s} | {:<20s}".format(
            shape.name, 
            str(shape.shape_type),
            type(shape).__name__,
            )
    if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
        shape_string += " | " + shape.ole_format.prog_id
    print(shape_string)
    if isinstance(shape, GroupShape):
        indent += 1
        for grp_shape in shape.shapes:
            print_shape_type(grp_shape, indent) 
    